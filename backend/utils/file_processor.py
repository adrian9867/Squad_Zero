import hashlib
import re
from io import BytesIO
from pathlib import Path
from typing import List

import PyPDF2
import docx
from fastapi import HTTPException, UploadFile

from config.config import settings
from utils.pdf_text_sanity import (
    looks_like_garbage,
    safe_strip_leaked_pdf_objects,
    ocr_pdf_bytes,        # shared OCR fallback (Tesseract -> vision LLM)
    ocr_image_bytes,      # shared image OCR fallback (Tesseract -> vision LLM)
    vision_ocr_image_bytes,  # direct vision OCR (used when caller already tried Tesseract)
)

# In-memory extraction cache keyed by SHA-256 of file bytes.
# Capped at _CACHE_MAX entries
from collections import OrderedDict as _OD
_CACHE_MAX = 100
_extraction_cache: _OD = _OD()


class FileProcessor:
    """
    Extracts plain text from all file types accepted by the frontend.

    Supported:
        Documents : pdf, doc, docx, txt, rtf
        Spreadsheets: xlsx, xls
        Presentations: ppt, pptx
        Images (OCR)  : jpg, jpeg, png, gif, webp, bmp, tiff
        E-books       : epub
    """

    # Every extension the frontend allows
    _ALL_SUPPORTED: set = {
        "pdf", "doc", "docx", "txt", "rtf",
        "xlsx", "xls",
        "ppt", "pptx",
        "jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "md", 
        "epub",
    }

    async def process_files(self, files: List[UploadFile]) -> str:
        """Process uploaded files and extract text, combining into one string.

        Results are cached by SHA-256 so repeated uploads of the same file
        skip extraction entirely.
        """
        combined_text = ""

        for file in files:
            ext = Path(file.filename or "").suffix.lstrip(".").strip().lower()

            if ext not in self._ALL_SUPPORTED:
                print(f"⚠️  Skipping {file.filename} (unsupported type: .{ext})")
                continue

            raw = await file.read()
            cache_key = hashlib.sha256(raw).hexdigest()

            if cache_key in _extraction_cache:
                print(f"📋 Cache hit: {file.filename} ({len(raw):,} bytes)")
                text = _extraction_cache[cache_key]
            else:
                text = self._extract(raw, ext, file.filename)
                # PDFs/DOCX/OCR can yield embedded NUL (0x00) bytes, which Postgres
                # text columns reject outright with "A string literal cannot contain
                # NUL (0x00) characters." Strip them right after extraction so the
                # bad byte never reaches the DB or the AI prompt.
                if text:
                    text = text.replace("\x00", "")
                if text:
                    _extraction_cache[cache_key] = text
                    if len(_extraction_cache) > _CACHE_MAX:
                        _extraction_cache.popitem(last=False)  # evict oldest entry
                    print(f"📥 Cached extraction: {file.filename} ({len(text):,} chars)")

            if text:
                combined_text += f"\n\n--- {file.filename} ---\n{text}"
            else:
                print(f"ℹ️  No text extracted from {file.filename}")

        if not combined_text.strip():
            raise HTTPException(
                status_code=400,
                detail=(
                    "No text could be extracted from the uploaded files. "
                    "For images, ensure they contain readable text. "
                    "For documents, ensure they are not password-protected."
                ),
            )

        return combined_text.strip()

    # Router

    # Magic-byte signatures for binary formats whose extraction requires a
    # dedicated parser (fitz for PDF, python-docx for DOCX, etc.). Used to
    # sniff the real format from the bytes when the filename extension is
    # wrong or missing — which happens when files are round-tripped through
    # the workspace file-content endpoint, where the frontend's
    # resolveDropFileName can rename a .pdf to .txt (e.g. when the file list
    # response includes a stale text column from a previous broken
    # extraction). Without this sniff, the .txt branch would just decode the
    # raw PDF bytes as UTF-8 (errors="ignore"), producing garbage that
    # looks_like_garbage() correctly flags — but too late, since the real
    # text was never extracted.
    _PDF_MAGIC = b"%PDF-"
    # DOCX/XLSX/PPTX are ZIP archives — PK\x03\x04
    _ZIP_MAGIC = b"PK\x03\x04"

    def _extract(self, content: bytes, ext: str, filename: str) -> str:
        """Route to the correct extractor based on file extension.

        Also sniffs magic bytes as a defense-in-depth measure: if the
        content starts with a known binary-format signature, route through
        the corresponding extractor regardless of the filename extension.
        This catches cases where the extension was lost or corrupted
        during a round-trip through the workspace file-content endpoint
        (the most common cause of the 'PDF encoding data' 400 error on
        folder-tree uploads).
        """
        try:
            # --- Magic-byte sniffing (defense in depth) ---
            # Only sniff when the extension is NOT already the correct
            # one — avoids redundant checks on the happy path.
            if ext != "pdf" and content[:5] == self._PDF_MAGIC:
                print(
                    f"⚠️  {filename} labeled .{ext} but starts with PDF magic "
                    f"(%PDF-) — extracting as PDF"
                )
                return self._extract_from_pdf(content)

            if ext not in ("docx", "xlsx", "pptx") and content[:4] == self._ZIP_MAGIC:
                # Could be DOCX, XLSX, or PPTX — try them in order. The
                # try/except inside each extractor handles wrong-format
                # cases gracefully (returns "" on failure).
                print(
                    f"⚠️  {filename} labeled .{ext} but starts with ZIP magic "
                    f"(PK) — trying DOCX/XLSX/PPTX extractors"
                )
                # Try DOCX first (most common), then XLSX, then PPTX
                for extractor_name in ("docx", "excel", "pptx"):
                    result = getattr(self, f"_extract_from_{extractor_name}")(content)
                    if result.strip():
                        return result
                # If none worked, fall through to the extension-based path

            # --- Extension-based routing (normal path) ---
            if ext == "pdf":
                return self._extract_from_pdf(content)
            elif ext in ("doc", "docx"):
                return self._extract_from_docx(content)
            elif ext in ("txt", "md"):
                return content.decode("utf-8", errors="ignore")
            elif ext == "rtf":
                return self._extract_from_rtf(content)
            elif ext in ("xlsx", "xls"):
                return self._extract_from_excel(content)
            elif ext in ("ppt", "pptx"):
                return self._extract_from_pptx(content)
            elif ext in ("jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff"):
                return self._extract_from_image(content)
            elif ext == "epub":
                return self._extract_from_epub(content)
            else:
                return ""
        except Exception as e:
            print(f"⚠️  Extraction failed for {filename}: {e}")
            return ""

    # Extractors

    def _looks_like_garbage(self, text: str) -> bool:
        """Does this look like leaked PDF structure rather than real text?
        See utils/pdf_text_sanity.py — shared across all PDF extractors."""
        return looks_like_garbage(text)

    def _extract_from_pdf(self, content: bytes) -> str:
        """Extract text from a PDF file.

        Tries PyMuPDF (fitz) first since it's already a dependency (used for
        OCR) and is far less prone to leaking raw PDF object-dictionary keys
        into the output than PyPDF2. Falls back to PyPDF2 if fitz isn't
        available. If the result from either looks like leaked PDF
        structure rather than real text, falls back to OCR, which sidesteps
        the parser entirely.
        """
        text = ""
        ocr_attempted = False

        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=content, filetype="pdf")
            page_texts = []
            for page in doc:
                try:
                    # sort=True restores natural reading order for PDFs whose
                    # content stream stores columns or glyphs out of order.
                    page_texts.append(page.get_text("text", sort=True) or "")
                except TypeError:
                    # Keeps the extractor compatible with light-weight test
                    # doubles and older PyMuPDF versions.
                    page_texts.append(page.get_text() or "")
            text = "\n".join(page_texts)
        except ImportError:
            print("ℹ️  PyMuPDF not installed — falling back to PyPDF2 for text extraction")
        except Exception as e:
            print(f"⚠️  PyMuPDF extraction failed, falling back to PyPDF2: {e}")

        # A malformed/oddly encoded PDF can return plenty of text that is
        # actually its internal object dictionary.  Try OCR first for that
        # case because it reads rendered page pixels rather than another
        # parser's interpretation of the broken text layer.
        if self._looks_like_garbage(text):
            print("⚠️  Extracted text looks like leaked PDF structure — falling back to OCR")
            ocr_attempted = True
            ocr_text = self._ocr_pdf(content)
            if ocr_text.strip():
                text = ocr_text

        # No selectable text at all (scanned PDF) → OCR
        if not text.strip():
            print("ℹ️  PDF has no selectable text — attempting OCR")
            return safe_strip_leaked_pdf_objects(self._ocr_pdf(content))

        # Selectable text exists but looks like leaked PDF internals → OCR.
        # This whole-document check runs on the raw extracted text, before
        # any surgical stripping below, since dilution from stripping could
        # otherwise mask a document that's genuinely all garbage.
        if self._looks_like_garbage(text) and not ocr_attempted:
            print("⚠️  Extracted text looks like leaked PDF structure — falling back to OCR")
            ocr_text = self._ocr_pdf(content)
            text = ocr_text if ocr_text.strip() else text

        # Unconditional final step: strip any leaked PDF object-dictionary
        # blocks that survived (or that OCR was never triggered for because
        # they were too small/diluted to trip the whole-document check).
        return safe_strip_leaked_pdf_objects(text)

    def _ocr_pdf(self, content: bytes) -> str:
        """OCR a scanned PDF by rendering each page as an image.

        Delegates to the shared ocr_pdf_bytes() in utils/pdf_text_sanity.py
        so this site picks up the same Tesseract-first / vision-LLM-fallback
        chain as every other PDF extractor (notes, flashcards, second brain).
        Previously this method had its own inline copy of the Tesseract call
        with no vision fallback, which meant scanned PDFs silently returned
        "" on any environment without the tesseract binary installed —
        producing 400 "No text could be extracted" errors on the quiz route.
        """
        return ocr_pdf_bytes(content)

    def _extract_from_docx(self, content: bytes) -> str:
        """Extract text from a DOCX file"""
        try:
            doc = docx.Document(BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            print(f"⚠️  DOCX extraction failed: {e}")
            return ""

    def _extract_from_rtf(self, content: bytes) -> str:
        """Extract text from an RTF file"""
        try:
            from striprtf.striprtf import rtf_to_text
            return rtf_to_text(content.decode("utf-8", errors="ignore"))
        except ImportError:
            print("⚠️  striprtf not installed. Run: pip install striprtf")
            # Fallback: strip RTF control words with regex
            import re
            raw = content.decode("utf-8", errors="ignore")
            text = re.sub(r"\{[^}]*\}", "", raw)
            text = re.sub(r"\\[a-z]+\d* ?", "", text)
            return text.strip()
        except Exception as e:
            print(f"⚠️  RTF extraction failed: {e}")
            return ""

    def _extract_from_excel(self, content: bytes) -> str:
        """Extract text from an XLSX/XLS file"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(BytesIO(content), read_only=True, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                text += f"\nSheet: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) for cell in row if cell is not None)
                    if row_text:
                        text += row_text + "\n"
            return text
        except ImportError:
            print("⚠️  openpyxl not installed. Run: pip install openpyxl")
            return ""
        except Exception as e:
            print(f"⚠️  Excel extraction failed: {e}")
            return ""

    def _extract_from_pptx(self, content: bytes) -> str:
        """Extract text from a PPTX/PPT file"""
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(content))
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\nSlide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text
        except ImportError:
            print("⚠️  python-pptx not installed. Run: pip install python-pptx")
            return ""
        except Exception as e:
            print(f"⚠️  PPTX extraction failed: {e}")
            return ""

    # --- Image MIME-type detection -------------------------------------
    # Used to label base64 image payloads sent to the vision LLM. Wrong
    # MIME types either get rejected by OpenRouter or cause the model to
    # mis-decode the image, so we sniff the magic bytes rather than
    # trusting the filename extension.
    _IMAGE_MAGIC = [
        (b"\x89PNG\r\n\x1a\n", "image/png"),
        (b"\xff\xd8\xff",      "image/jpeg"),
        (b"GIF87a",            "image/gif"),
        (b"GIF89a",            "image/gif"),
        (b"BM",                "image/bmp"),
        (b"II*\x00",           "image/tiff"),
        (b"MM\x00*",           "image/tiff"),
        # WebP: bytes 0-3 = "RIFF", bytes 8-11 = "WEBP"
    ]

    @classmethod
    def _guess_image_mime_type(cls, content: bytes) -> str:
        """Return the image MIME type from magic bytes, defaulting to png."""
        if not content:
            return "image/png"
        # WebP needs a split check
        if len(content) >= 12 and content[:4] == b"RIFF" and content[8:12] == b"WEBP":
            return "image/webp"
        for magic, mime in cls._IMAGE_MAGIC:
            if content.startswith(magic):
                return mime
        return "image/png"

    def _extract_from_image(self, content: bytes) -> str:
        """Extract text from an image.

        Strategy:
          1. Try local Tesseract OCR first, with light preprocessing
             (flatten transparency, upscale small images, grayscale +
             threshold) and two PSM modes. This maximises Tesseract
             accuracy on screenshots and exported lecture slides, which
             commonly have small text and anti-aliased edges.
          2. If Tesseract is unavailable, returns nothing usable, or
             raises any error (e.g. corrupt image), fall back to vision-
             LLM OCR via OpenRouter. The vision LLM doesn't need any
             system binary — only OPENROUTER_API_KEY in env — and is
             significantly more accurate on small text, code, formulas,
             and coloured backgrounds.

        The fallback is what makes image uploads work on environments
        where the `tesseract` binary isn't installed (common in
        containerised/managed hosts). Without it, every image upload
        silently returned "" and produced a 400 "No text could be
        extracted" error on the quiz route.
        """
        mime_type = self._guess_image_mime_type(content)

        # --- Path 1: Tesseract with preprocessing ---
        tesseract_text = self._tesseract_image_ocr(content)
        if tesseract_text.strip():
            return tesseract_text

        # --- Path 2: vision-LLM fallback (no system binary needed) ---
        try:
            vision_text = vision_ocr_image_bytes(content, mime_type)
            if vision_text.strip():
                return vision_text
        except Exception as e:
            print(f"⚠️  Vision OCR fallback failed for image: {e}")

        # Neither path yielded anything — return whatever Tesseract gave
        # us (possibly "") so the caller's existing empty-result handling
        # kicks in uniformly.
        return tesseract_text

    def _tesseract_image_ocr(self, content: bytes) -> str:
        """Run Tesseract OCR with light preprocessing. Returns "" if
        Tesseract is unavailable, returns nothing usable, or raises any
        error. Separated out so _extract_from_image can call vision OCR
        only when this path genuinely produced nothing."""
        try:
            import shutil
            if shutil.which("tesseract") is None:
                # Don't even print here — this is the expected case on
                # systems without tesseract installed, and the vision
                # fallback below will log what it's doing. Printing here
                # produced noise on every single image upload.
                return ""

            import pytesseract
            from PIL import Image, ImageFilter, ImageOps

            img = Image.open(BytesIO(content))

            # Flatten transparency onto white instead of allowing transparent
            # pixels to become black text-shaped noise.
            if "A" in img.getbands():
                background = Image.new("RGBA", img.size, "white")
                img = Image.alpha_composite(background, img.convert("RGBA")).convert("RGB")
            elif img.mode not in ("RGB", "L"):
                img = img.convert("RGB")

            # Upscaling helps Tesseract resolve the small fonts used in
            # screenshots.  Keep large uploads bounded to avoid huge OCR jobs.
            max_dimension = max(img.size)
            if max_dimension < 1800:
                scale = min(3.0, 1800 / max_dimension)
                img = img.resize(
                    (round(img.width * scale), round(img.height * scale)),
                    Image.Resampling.LANCZOS,
                )
            elif max_dimension > 3200:
                scale = 3200 / max_dimension
                img = img.resize(
                    (round(img.width * scale), round(img.height * scale)),
                    Image.Resampling.LANCZOS,
                )

            gray = ImageOps.grayscale(img)
            gray = ImageOps.autocontrast(gray)
            sharp = gray.filter(ImageFilter.SHARPEN)
            threshold = sharp.point(lambda pixel: 255 if pixel > 175 else 0)
            variants = (img, sharp, threshold)

            best_text = ""
            best_score = -1
            for variant in variants:
                for psm in (6, 11):
                    candidate = pytesseract.image_to_string(
                        variant,
                        config=f"--oem 3 --psm {psm}",
                    ).strip()
                    # Prefer text with real words, while retaining punctuation
                    # and numbers that may matter in equations/code.
                    score = len(re.findall(r"[A-Za-z]{2,}", candidate)) * 10 + len(candidate)
                    if score > best_score:
                        best_text, best_score = candidate, score
            return best_text
        except ImportError:
            print("⚠️  pytesseract or Pillow not installed — falling back to vision OCR for image")
            return ""
        except Exception as e:
            # This includes a corrupt image bytes case.  Vision OCR below
            # is the last-resort fallback for genuinely unexpected errors.
            print(f"⚠️  Tesseract image OCR failed: {e} — trying vision OCR fallback")
            return ""

    def _extract_from_epub(self, content: bytes) -> str:
        """Extract text from an EPUB file"""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            # write to a temp file as ebooklib requires file path 
            import os
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".epub", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name

            try:
                book = epub.read_epub(tmp_path)
                text = ""
                for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text += soup.get_text(separator="\n") + "\n"
                return text.strip()
            finally:
                os.unlink(tmp_path)

        except ImportError:
            print(
                "⚠️  ebooklib or beautifulsoup4 not installed.\n"
                "    Run: pip install ebooklib beautifulsoup4"
            )
            return ""
        except Exception as e:
            print(f"⚠️  EPUB extraction failed: {e}")
            return ""
